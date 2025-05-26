import time
import os
import re
import threading
import tempfile
import requests
import pyautogui
import pyperclip
import speech_recognition as sr
from gtts import gTTS
import pygame
from selenium import webdriver
from selenium.webdriver.edge.service import Service
from pptx import Presentation
from pptx.util import Inches
from fuzzywuzzy import process



# Initialize pygame mixer for audio playback
pygame.mixer.init()

# Configurations
GEMINI_API_KEY = "AIzaSyBAn0xxX4IRnP_gOH4P6MUMakw6PPd4p_Q"  
PEXELS_API_KEY = "Ee5M2q5YipMtBcTgfRDFyo2rgIVzdDce7xbTtcaPrBky4DgjGaiNwcyY"

edge_service = Service("C:/WebDriver/msedgedriver.exe")
driver = None  # Lazy initialization for the browser

# Global flag for voice mode
listening = False

def shorten_text(text):
    """Return only the first sentence for brevity."""
    if '.' in text:
        return text.split('.')[0].strip() + '.'
    return text.strip()

def speak(text):
    """
    Convert text to speech using gTTS and play it with pygame.
    Only a short version (the first sentence) is spoken.
    After playback, a background thread removes the temporary MP3.
    """
    short_text = shorten_text(text)
    print(f"🤖 AI: {short_text}")

    try:
        tts = gTTS(short_text, lang="en")
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            tmp_path = fp.name
        tts.save(tmp_path)
        pygame.mixer.music.load(tmp_path)
        pygame.mixer.music.play()
        while pygame.mixer.music.get_busy():
            pygame.time.Clock().tick(10)
    except Exception as e:
        print(f"Speech Error: {e}")
    finally:
        def delayed_remove(path):
            time.sleep(2)
            try:
                if hasattr(pygame.mixer.music, "unload"):
                    pygame.mixer.music.unload()
            except Exception as unload_error:
                print(f"Error during unload: {unload_error}")
            for attempt in range(3):
                try:
                    os.remove(path)
                    break
                except Exception as ex:
                    print(f"Attempt {attempt+1}: Error removing temporary file {path}: {ex}")
                    time.sleep(1)
        threading.Thread(target=delayed_remove, args=(tmp_path,), daemon=True).start()

def listen():
    """
    Listen for a voice command using the microphone.
    Adjust for ambient noise and return the recognized text in lowercase.
    """
    r = sr.Recognizer()
    try:
        with sr.Microphone() as source:
            r.adjust_for_ambient_noise(source)
            print("🎤 Listening...")
            audio = r.listen(source, timeout=5, phrase_time_limit=5)
        command = r.recognize_google(audio, language="en-US").lower().strip()
        print(f"🎙 You said: {command}")
        return command
    except Exception as e:
        print(f"❌ Listen error: {e}")
        return ""

def search_web(query):
    """
    Open a Google search in Microsoft Edge for the provided query.
    """
    global driver
    print(f"🌍 Searching for: {query}...")
    speak(f"Searching for {query}")
    try:
        if driver is None:
            driver = webdriver.Edge(service=edge_service)
        search_url = f"https://www.google.com/search?q={query.replace(' ', '+')}"
        driver.get(search_url)
    except Exception as e:
        print(f"❌ Browser Error: {e}")
        speak("Couldn't open the browser.")

def clean_text(text):
    """Clean text by removing unwanted characters."""
    if not text:
        return "No answer."
    return re.sub(r"[^a-zA-Z0-9,.!? \n]", "", text).strip()

def get_ai_generated_text(prompt, retries=3, use_clean=True):
    """
    Use Gemini's text-generation endpoint to get AI‑generated text.
    If use_clean is True the returned text is cleaned.
    """
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"
    headers = {"Content-Type": "application/json"}
    data = {"contents": [{"parts": [{"text": prompt}]}]}
    for _ in range(retries):
        try:
            response = requests.post(url, json=data, headers=headers, timeout=15)
            if response.status_code == 200:
                candidates = response.json().get("candidates", [])
                if candidates and "content" in candidates[0]:
                    parts = candidates[0]["content"].get("parts", [])
                    if parts and "text" in parts[0]:
                        raw_text = parts[0]["text"]
                        return clean_text(raw_text) if use_clean else raw_text
            print(f"❌ API Error: {response.status_code}")
        except Exception as e:
            print(f"❌ API Request Error: {e}")
    return "Sorry, I couldn't get a response from the AI."

def open_and_paste(prompt, disable_cleaning=False):
    """
    Generate content via Gemini and write it inside a new blank document in Microsoft Word.
    Instead of manually clicking, it just presses 'Enter' after opening Word.
    """
    print(f"📝 Generating content on '{prompt}'...")
    speak(f"Generating content on {prompt}")
    content = get_ai_generated_text(prompt, use_clean=(not disable_cleaning))
    formatted = f"**{prompt.title()}**\n\n{content}"
    pyperclip.copy(formatted)

    # Step 1: Open Windows Search
    pyautogui.press("win")
    time.sleep(1.5)  # Small delay for Windows Search to activate

    # Step 2: Type "Microsoft Word" and open it
    pyautogui.write("Microsoft Word")  # Full name instead of "Word"
    time.sleep(1.5)  # Allow Windows Search time to process input

    pyautogui.press("enter")  # Open Word
    time.sleep(4)  # Increased wait time for Word to load

    # Step 3: Automatically start a blank document (just press Enter)
    pyautogui.press("enter")
    time.sleep(2)  # Allow Word to switch to blank document

    # Step 4: Paste content inside the blank Word document
    try:
        pyautogui.hotkey("ctrl", "v")  # Paste content
        speak("Done writing in Microsoft Word.")
    except Exception as e:
        speak(f"Error writing: {str(e)}")



def open_application(app_name):
    """
    Open an application by simulating keypresses (using Windows search).
    """
    print(f"🖥 Opening {app_name}...")
    speak(f"Opening {app_name}")
    pyautogui.hotkey("win", "s")
    time.sleep(1)
    pyautogui.typewrite(app_name)
    time.sleep(1)
    pyautogui.press("enter")

def generate_image(prompt):
    """
    Generate an image using Gemini's (hypothetical) image-generation endpoint.
    Returns the path to a temporary image file.
    """
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateImage?key={GEMINI_API_KEY}"
    headers = {"Content-Type": "application/json"}
    data = {"prompt": prompt}

    try:
        response = requests.post(url, json=data, headers=headers, timeout=30)
        
        # Check for API success
        if response.status_code != 200:
            print(f"Error: API responded with {response.status_code}")
            return None

        json_data = response.json()
        
        # Validate presence of image URL
        image_url = json_data.get("imageUrl")
        if not image_url:
            print("Error: No image URL found in response.")
            return None
        
        # Download image content
        img_response = requests.get(image_url)
        if img_response.status_code != 200:
            print(f"Error downloading image, HTTP Status: {img_response.status_code}")
            return None
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_img_file:
            temp_img_file.write(img_response.content)
            return temp_img_file.name

    except requests.exceptions.RequestException as e:
        print(f"Network error while generating image: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")

    return None

def display_image(image_path):
    """
    Display the generated image using the default image viewer.
    """
    if not image_path or not os.path.exists(image_path):
        print("Error: Image file not found.")
        return
    
    print("Displaying generated image...")
    os.startfile(image_path)

#1ppt


def fetch_image(topic):
    """Search for a relevant image using Pexels API."""
    url = f"https://api.pexels.com/v1/search?query={topic}&per_page=1"
    headers = {"Authorization": PEXELS_API_KEY}
    
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        data = response.json()
        return data["photos"][0]["src"]["original"]  # Get first image URL
    else:
        print(f"❌ Error fetching image: {response.status_code}")
        return None

PEXELS_API_KEY = "your_pexels_api_key"

def fetch_images(topic, num_images):
    """Search for multiple relevant images using Pexels API."""
    url = f"https://api.pexels.com/v1/search?query={topic}&per_page={num_images}"
    headers = {"Authorization": PEXELS_API_KEY}
    
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        data = response.json()
        return [photo["src"]["original"] for photo in data["photos"]]  # List of image URLs
    else:
        print(f"❌ Error fetching images: {response.status_code}")
        return []  # Return empty list if API fails

import requests
from pptx import Presentation
from pptx.util import Inches
import os

PEXELS_API_KEY = "your_pexels_api_key"

def fetch_images(topic, num_images):
    """Search for relevant images using Pexels API."""
    url = f"https://api.pexels.com/v1/search?query={topic}&per_page={num_images}"
    headers = {"Authorization": PEXELS_API_KEY}
    
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        data = response.json()
        return [photo["src"]["original"] for photo in data["photos"]]  # List of image URLs
    else:
        print(f"❌ Error fetching images: {response.status_code}")
        return []  # Return empty list if API fails

def create_presentation():
    """Generates a structured PowerPoint presentation with topic-based images."""
    user_input = input("Enter 'pptx' to create a presentation or any other text to exit: ").strip().lower()
    
    if user_input != "pptx":
        print("Exited presentation mode.")
        return

    topic = input("Enter the presentation topic: ").strip()
    print(f"Creating PPTX on topic: {topic}")
    
    prompt = (
        f"Generate a well-organized PowerPoint presentation on {topic}. "
        "Ensure structured sections:\n"
        "- **Slide 1**: Title & Introduction\n"
        "- **Slide 2-4**: Facts, Analysis, Explanation (Paragraphs, not just bullet points)\n"
        "- **Slide 5**: Conclusion (Summarize key findings and provide insights)\n"
        "Do NOT generate vague bullet points. Provide detailed, **fact-based** explanations."
    )

    presentation_text = get_ai_generated_text(prompt)
    slides_content = [slide.strip() for slide in presentation_text.split('---') if slide.strip()]

    if not slides_content:
        print("❌ Could not generate structured slide content.")
        return

    prs = Presentation()
    image_urls = fetch_images(topic, num_images=len(slides_content))  # Get same number of images as slides

    for index, slide_text in enumerate(slides_content[:5]):  # Limit slides to 5
        lines = slide_text.splitlines()
        if not lines:
            continue

        title = lines[0]
        content = "\n".join(lines[1:]) if len(lines) > 1 else "No additional content."

        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout
        slide.shapes.title.text = title
        content_box = slide.shapes.placeholders[1]
        content_box.text = content  # Full sentences, not messy bullet points

        # Add images sequentially based on the slide index
        if index < len(image_urls):
            image_url = image_urls[index]
            try:
                img_path = os.path.join(os.getcwd(), f"image_{index}.jpg")
                img_data = requests.get(image_url).content
                with open(img_path, "wb") as img_file:
                    img_file.write(img_data)

                slide.shapes.add_picture(img_path, Inches(0.5), Inches(3), width=Inches(5), height=Inches(3))
                os.remove(img_path)
            except Exception as e:
                print(f"⚠️ Error adding image to slide {index + 1}: {e}")

    ppt_path = os.path.join(os.getcwd(), "presentation.pptx")
    prs.save(ppt_path)
    
    print(f"✅ Structured presentation saved at {ppt_path}")
    speak("Your PowerPoint presentation is ready.")

    if os.path.exists(ppt_path):
        os.startfile(ppt_path)  # Open PowerPoint only if the file exists
    else:
        print("⚠️ Error: Presentation file not found.")
        speak("I couldn't find the PowerPoint file to open.")



# AI-Generated Image Handling
def generate_image(prompt):
    """
    Generate an image using Gemini's (hypothetical) image-generation endpoint.
    Returns the path to a temporary image file.
    """
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateImage?key={GEMINI_API_KEY}"
    headers = {"Content-Type": "application/json"}
    data = {"prompt": prompt}

    try:
        response = requests.post(url, json=data, headers=headers, timeout=30)

        if response.status_code != 200:
            print(f"Error: API responded with {response.status_code}")
            return None

        json_data = response.json()
        image_url = json_data.get("imageUrl")

        if not image_url:
            print("Error: No image URL found in response.")
            return None
        
        img_response = requests.get(image_url)
        if img_response.status_code != 200:
            print(f"Error downloading image, HTTP Status: {img_response.status_code}")
            return None
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_img_file:
            temp_img_file.write(img_response.content)
            return temp_img_file.name

    except requests.exceptions.RequestException as e:
        print(f"Network error while generating image: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")

    return None

# Word Processing Functionality
def open_and_paste(prompt, disable_cleaning=False):
    """
    Generate content via Gemini and write it inside a new blank document in Microsoft Word.
    Instead of pasting anywhere, it ensures Microsoft Word opens and a blank document is ready.
    """
    print(f"📝 Generating content on '{prompt}'...")
    speak(f"Generating content on {prompt}")
    content = get_ai_generated_text(prompt, use_clean=(not disable_cleaning))
    formatted = f"**{prompt.title()}**\n\n{content}"
    pyperclip.copy(formatted)

    # Open Microsoft Word
    pyautogui.press("win")
    time.sleep(1.5)
    pyautogui.write("Microsoft Word")
    time.sleep(1.5)
    pyautogui.press("enter")
    time.sleep(4)
    pyautogui.press("enter")
    time.sleep(2)

    try:
        pyautogui.hotkey("ctrl", "v")  # Paste content
        speak("Done writing in Microsoft Word.")
    except Exception as e:
        speak(f"Error writing: {str(e)}")

# Command Routing
def process_command(cmd):
    """Route the command to the appropriate functionality."""
    global listening
    if not cmd:
        return

    if cmd.startswith("write "):
        write_prompt = cmd.replace("write", "").strip()

        if "code" in write_prompt.lower():
            open_and_paste(write_prompt, disable_cleaning=True)
        else:
            open_and_paste(write_prompt, disable_cleaning=False)

    elif cmd.startswith("what is "):
        response = get_ai_generated_text(cmd.replace("what is", "").strip())
        speak(response)

    elif cmd.startswith("open "):
        open_application(cmd.replace("open", "").strip())

    elif cmd.startswith("search "):
        search_web(cmd.replace("search", "").strip())

    elif cmd == "pptx":  # 🔥 New trigger for PowerPoint generation
        create_presentation()

    elif cmd == "start":
        if not listening:
            listening = True
            threading.Thread(target=voice_loop, daemon=True).start()
            print("🟢 Voice mode started.")
            speak("Voice mode activated.")
        else:
            speak("Voice mode is already active.")

    elif cmd == "stop":
        listening = False
        speak("Voice mode deactivated.")

    elif cmd in ["exit", "quit"]:
        if driver:
            driver.quit()
        exit()

    else:
        response = get_ai_generated_text(cmd)
        speak(response)

# Continuous Voice Processing
def voice_loop():
    """Process voice commands continuously when voice mode is active."""
    global listening
    MIN_LOOP_INTERVAL = 5  # seconds per iteration
    while listening:
        start = time.time()
        cmd = listen()
        if cmd:
            process_command(cmd)
        elapsed = time.time() - start
        if elapsed < MIN_LOOP_INTERVAL:
            extra = MIN_LOOP_INTERVAL - elapsed
            speak("It's too fast.")
            time.sleep(extra)





def main():
    """Main assistant loop. Accept text input if voice mode is not active."""
    global listening
    print("🤖 AI Assistant ready!")
    print("Type 'start' for voice mode or enter commands manually:")
    while True:
        if not listening:
            cmd = input("💬: ").strip().lower()
            process_command(cmd)
        else:
            time.sleep(1)

if __name__ == "__main__":
    main()